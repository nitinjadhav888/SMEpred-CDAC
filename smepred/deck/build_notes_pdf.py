"""
build_notes_pdf.py — make a "Notes Pages" handout PDF.

PowerPoint COM's ExportAsFixedFormat with OutputType=NotesPages is brittle from
PowerShell/Python late-binding. So we build the handout ourselves: each page shows the
slide thumbnail (rendered via PowerPoint to img/Slide{N}.PNG earlier) plus the speaker
notes pulled straight from the .pptx file. The output mirrors the standard PowerPoint
Notes Pages handout but is produced deterministically by ReportLab.
"""
from pathlib import Path
from pptx import Presentation
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Image, Paragraph, Spacer, PageBreak, Table, TableStyle
)
from reportlab.lib.colors import HexColor
from reportlab.lib.units import inch

DECK_DIR = Path(__file__).parent
PPTX = DECK_DIR / "SMEpred_Pitch_Deck.pptx"
IMG  = DECK_DIR / "img"
OUT  = DECK_DIR / "SMEpred_Pitch_Deck_with_Notes.pdf"

PRIMARY = HexColor("#0B2239")
ACCENT  = HexColor("#02C39A")
MUTED   = HexColor("#64748B")

styles = getSampleStyleSheet()
title_style = ParagraphStyle(
    "title", parent=styles["Heading2"], fontSize=14, leading=18,
    textColor=PRIMARY, spaceAfter=4)
kicker = ParagraphStyle(
    "kicker", parent=styles["Normal"], fontSize=8.5, leading=11,
    textColor=ACCENT, spaceAfter=2)
note_head = ParagraphStyle(
    "nh", parent=styles["Normal"], fontSize=9.5, leading=12,
    textColor=MUTED, spaceAfter=4)
note_body = ParagraphStyle(
    "nb", parent=styles["Normal"], fontSize=10.5, leading=14, textColor=PRIMARY,
    spaceAfter=4, allowOrphans=1)


def slide_thumbnail_path(n: int) -> Path:
    return IMG / f"Slide{n}.PNG"


def extract_notes() -> list[str]:
    prs = Presentation(str(PPTX))
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes.append(slide.notes_slide.notes_text_frame.text.strip())
        else:
            notes.append("")
    return notes


def build():
    notes = extract_notes()
    total = len(notes)
    doc = SimpleDocTemplate(
        str(OUT), pagesize=LETTER,
        leftMargin=0.55*inch, rightMargin=0.55*inch,
        topMargin=0.55*inch, bottomMargin=0.55*inch,
        title="SMEpred Pitch — Speaker Notes Handout",
    )
    story = []
    # cover
    story.append(Paragraph("SMEpred — Speaker Notes Handout", title_style))
    story.append(Paragraph("Pitch deck for C-DAC Pune internship presentation", kicker))
    story.append(Spacer(1, 0.15*inch))
    story.append(Paragraph(
        "Each page below shows one slide alongside its speaker notes. "
        "Print the slides-only PDF for the audience and this handout for yourself.",
        note_body))
    story.append(PageBreak())

    img_w = 5.6 * inch              # slide thumbnail width
    img_h = img_w * 9 / 16          # 16:9 aspect

    for i, note in enumerate(notes, start=1):
        thumb = slide_thumbnail_path(i)
        story.append(Paragraph(f"Slide {i} of {total}", kicker))
        if thumb.exists():
            img = Image(str(thumb), width=img_w, height=img_h)
            img.hAlign = "CENTER"
            story.append(img)
        else:
            story.append(Paragraph(f"(thumbnail not found: {thumb.name})", note_head))
        story.append(Spacer(1, 0.18*inch))
        story.append(Paragraph("Speaker notes", note_head))
        # split by blank lines into paragraphs; preserve simple line breaks
        chunks = [c.strip() for c in note.split("\n\n") if c.strip()]
        for c in chunks:
            text = c.replace("\n", "<br/>")
            story.append(Paragraph(text, note_body))
        if i < total:
            story.append(PageBreak())

    doc.build(story)
    print(f"Wrote {OUT}  ({OUT.stat().st_size/1024:.0f} KB)  with {total} slides + notes.")


if __name__ == "__main__":
    build()
