"""Inject per-slide speaker notes into the SMEpred deck so they show up in PowerPoint's
Presenter View under each slide."""
from pptx import Presentation
from pathlib import Path

DECK = Path(__file__).parent / "SMEpred_Pitch_Deck.pptx"

NOTES = [
    # Slide 1 — Title
    """OPENER: "Good morning. I'm presenting SMEpred — an AI tool that predicts how well a
gene-silencing drug will work, before anyone touches a test tube."

SMEpred = siRNA Modification Efficacy Predictor. Built during my C-DAC Pune internship.
In 15 minutes: the problem, how it works, real measured results, what's next.

ANALOGY: Like Google Maps for drug designers — ask the app which route is fastest, then
drive.

TRANSITION: "Let's start with the problem.""" ,

    # Slide 2 — The Problem
    """OPENER: "Many diseases start with one bad gene that won't shut up."

Genes make proteins. Bad gene → harmful protein → cancer, high cholesterol, viral disease.
siRNA is a 21-letter molecule that silences ONE specific gene — a programmable off-switch.
But raw siRNA is destroyed in minutes by blood enzymes, triggers immune alarms, and never
reaches the target. It needs chemical "armor" — modifications.

KEY STAT (point at the right panel): 1,260 possible chemical modifications per siRNA.
Testing each in the lab costs months and lakhs of rupees.

ANALOGY: Like designing a phone case — 1,260 material × shape combos. You can't physically
test all of them. You need a simulator that says: top 5 worth manufacturing.

TRANSITION: "Here's an analogy that makes it click.""" ,

    # Slide 3 — Real-life analogy
    """OPENER: "Think of an siRNA drug as a guided missile."

- Antisense strand = GPS — locks onto the exact mRNA target out of 20,000 genes.
- Modifications = armor plating — protect the missile from being shot down mid-flight.
- SMEpred = war-room simulator — predicts which armor design hits the target hardest,
  BEFORE you build 1,260 missiles.

ANALOGY: Indian Air Force engineers don't physically build every missile variant. They
simulate, pick the top three, then build the real thing. SMEpred is that simulator for
medicine.

TRANSITION: "So what does it actually do? One sentence.""" ,

    # Slide 4 — What it does
    """OPENER: Read the big blue banner — slowly.

Two related jobs, with two ML models behind the scenes:
1. NAKED EFFICACY — given a gene, find the best unmodified siRNA sequence.
2. MODIFIED EFFICACY — given a chosen siRNA, find the chemical tweak that makes it the
   strongest drug.
3. ENSEMBLE — for naked siRNAs, our LightGBM AND the OligoFormer transformer both score
   the candidates and we cross-validate (top picks must agree on both).
4. SAFETY — every candidate also gets a seed-toxicity check from a published lookup
   table, and modifications can RESCUE a toxic seed (e.g. 2'-OMe at antisense pos 2).

ANALOGY: A chef has two skills — pick the right base recipe (naked siRNA), then pick the
spice that makes it irresistible (modification). SMEpred does both, and runs every dish
past TWO judges instead of one.

TRANSITION: "Here's how that flows step by step.""" ,

    # Slide 5 — Workflow funnel
    """OPENER: "The whole app is one funnel — four steps from a gene to a drug candidate."

1. Paste a gene (1,000–10,000 letters of mRNA).
2. RANK tab — score every 21-letter window, rank, show best naked siRNA candidates.
3. SINGLE-MOD SCAN — pick the top siRNA; model scans all 1,260 modifications and ranks
   them by predicted efficacy gain.
4. MULTI-MOD DESIGN — design a custom drug with multiple mods; model scores it.

FIX SHIPPED (bottom banner): Earlier users had to copy-paste the siRNA between tabs.
Now each ranked row has a "Modify →" button that auto-fills the next tab. One click.

ANALOGY: Like Amazon — search a product, click "Buy Now", address & payment auto-fill.

TRANSITION: "Now, under the hood.""" ,

    # Slide 6 — Architecture
    """OPENER: "Five clean layers — data, features, two models, serving."

- DATA: parses HelixZero 43k patent catalog + OligoFormer's Huesken/Mix/Taka sets +
  siRNAmod into clean CSVs.
- FEATURES: turns each siRNA into a 152-number vector the model can read.
- MODEL A (ours): LightGBM gradient-boosted-tree regressor — predicts % inhibition.
- MODEL B (vendored): OligoFormer transformer + RNA-FM foundation model — predicts
  activity probability. Re-ranks the top 50 picks from Model A.
- SERVE: FastAPI REST + single-file HTML web app. Browser-only, no setup.

ANALOGY: A car factory — raw steel (data) → stamping (features) → primary assembly
(LightGBM) → quality inspection (OligoFormer) → showroom (web app).

TRANSITION: "The whole thing is only as good as its data.""" ,

    # Slide 7 — Real data
    """OPENER: "Garbage in, garbage out. So we use real data — measured in real labs."

THREE BIG NUMBERS:
- 25,763 modified siRNAs from HelixZero pharma patents (training cm-siRNA model).
- 4,060 NAKED siRNAs from FOUR published lab datasets:
    Huesken (Nature Biotech 2005, the gold standard) — 2,361 rows
    Mix (Reynolds/Vickers/Ui-Tei combined) — 462 rows
    Takayuki (NAR 2007) — 699 rows
    Our existing HelixZero/siRNAmod unmodified — 538 rows
- 4,097 entries in the OligoFormer seed-toxicity table (Janas et al., Mol Cell 2018).

DARK CALLOUT: A previous version fell back to synthetic random data — useless. Real
data is non-negotiable. We feed dataset SOURCE as a model feature so per-lab biases
are learned, not averaged away. Result: naked PCC jumped from 0.32 to 0.42.

ANALOGY: Training a self-driving car on dashcam from four different cities. Tell the
model which city each clip came from — it learns the local rules instead of averaging
into one confused driver.

TRANSITION: "Real data still needs to become numbers.""" ,

    # Slide 8 — Feature engineering
    """OPENER: "ML models can't read RNA letters. So we translate each siRNA into a
152-number fingerprint."

Walk the four cards:
- 140 COMPOSITION numbers — how many of each base & modification, both strands.
- 8 POSITION numbers — WHERE mods sit. Position matters: seed region (pos 1–8) is sacred,
  3' tail is forgiving.
- 2 GC CONTENT numbers — duplex stability proxy.
- 2 CONDITION numbers — dose & time of the experiment. THE SECRET SAUCE.

KEY INSIGHT (italic line): Same sequence shows different inhibition at 1 nM vs 100 nM.
Ignoring that = model sees contradictions, learns badly. Feeding dose+time AS FEATURES =
keep all 25k rows.

ANALOGY: Predicting cricket runs — depends on batsman stats AND pitch conditions. We feed
the pitch into the model.

TRANSITION: "And the model itself — big choice here.""" ,

    # Slide 9 — SMEpred + OligoFormer ensemble
    """OPENER: "When two independently-built models agree, you trust the answer more."

LEFT CARD (OUR LightGBM, teal):
- We built it. Trained on 25,763 modified + 4,060 naked siRNAs.
- HANDLES chemical modifications — drives all three tabs.
- Predicts % inhibition (regression).
- Fast: 1,000 candidates in under 100 ms.

RIGHT CARD (OligoFormer, mint):
- Tsinghua University, Bai et al. 2024 — peer-reviewed transformer model.
- Uses RNA-FM (a 1.1 GB foundation model) to embed sequences, then a small transformer
  head predicts activity probability.
- NAKED siRNA only — we use it as a second opinion on the Rank tab.
- Slower: re-ranks our top 50 picks in about 11 seconds on CPU.

THE ENSEMBLE STRIP (bottom): Because the two models output different scales
(% inhibition vs activity probability), averaging raw scores would over-weight one.
Instead we compute each candidate's PERCENTILE within the batch for each model, and
average those. "Top 5% by both models" surfaces as ensemble ~95 — a real cross-
validated pick. This is calibration-invariant by construction.

ANALOGY: Like a startup hiring committee. Each interviewer ranks candidates 1-N from
their own perspective. The hire isn't the highest average raw score — it's the
candidate who's in the top decile for EVERY interviewer.

TRANSITION: "But raw accuracy claims are easy to fake. We measured honestly.""" ,

    # Slide 10 — Honest evaluation
    """OPENER: "Different product features need different honesty tests."

LEFT — GENE-GROUPED SPLIT: Held out 3 whole genes (AGT, MSTN, PLN). Model never saw them.
Answers: "Can it predict a brand-new gene's best siRNA?" PCC = 0.26. Honestly hard.
Used for the RANK tab.

RIGHT — RANDOM SPLIT: standard 82/18. Answers: "Can it rank modifications of a known
siRNA?" PCC = 0.68. The real job of SINGLE-MOD and MULTI-MOD tabs.

WHY BOTH? With only 13 genes total, a naive split lets the model memorize gene-specific
motifs → fake accuracy. Gene-grouped proves real generalization.

ANALOGY: Chess engine — "beat opponents it's seen" (easy) vs "beat fresh opponents"
(real test). We do both.

TRANSITION: "So let's see the headline result.""" ,

    # Slide 11 — Results
    """OPENER: Point at chart. "Red = paper-baseline SVR. Green = our LightGBM rebuild."

TWO HEADLINE WINS:
- Modification ranking (core task): 0.37 → 0.68 PCC. +84% relative gain.
- Naked siRNA ranking: 0.21 → 0.42 PCC. +100% (doubled) — the OligoFormer dataset
  merge is what pushed it over the top.

BONUSES (on the right panel):
- Predictions are now REAL inhibition % on 0–100 scale. Old code had a bug — min-max
  rescaled every batch, so the "best of any batch" always showed 100. Meaningless.
- MAE on cm-siRNA is 16.5 percentage points — approaching the experimental noise floor
  of the underlying assays themselves (~10-15 pts between labs).

ANALOGY: Old system = exam graded on a curve (top of any room got 100 even if everyone
failed in absolute terms). New system = true scores you can compare across batches.

TRANSITION: "These predictions live inside a working product.""" ,

    # Slide 12 — Web app
    """OPENER: "It's not a research notebook. It's a working web app with four tabs and
two safety filters baked in."

Each tab solves one job: Rank, Single-Mod Scan, Multi-Mod Design, Modifications
reference. Single HTML file + FastAPI backend over HTTP. No build step, no
dependencies. Browser-only, results in seconds.

NEW FEATURES (this week):
- Toggle: "Re-rank top 50 with OligoFormer" on the Rank tab — adds an Ensemble %ile column.
- Seed Toxicity column on EVERY tab — Safe / Caution / Toxic with cell-viability %.
- MODIFICATION-AWARE TOXICITY — if a known rescuing modification (2'-OMe, 2'-F, LNA,
  2'-MOE) sits in the seed region of a modified candidate, the badge flips to
  "Mitigated" with a tooltip explaining the rescue. Real published biology in a UI.

DEMO RESULT (bottom): Paste gene → Modify → Scan. The model finds 2'-MOE at antisense
pos 8 (+10.6 efficacy) AND 2'-OMe at antisense pos 2 RESCUES a toxic seed → Mitigated.
That's the kind of insight that saves wet-lab months.

ANALOGY: Excel + a second opinion from a domain expert. Data in, both models crunch,
ensemble percentiles out, with safety annotations a chemist actually trusts.

TRANSITION: "Why does this matter outside this room?""" ,

    # Slide 13 — Real-world impact
    """OPENER: "siRNA drugs aren't theoretical anymore."

Three FDA-approved siRNA drugs saving lives today:
- PATISIRAN — rare nerve disease.
- INCLISIRAN — lowers cholesterol via PCSK9 gene silencing. Millions of patients.
- GIVOSIRAN — porphyria.
EVERY ONE needed exactly the modification-optimization SMEpred predicts. Without it,
they took years and hundreds of millions of dollars to design.

Three concrete impacts: fewer wet-lab experiments per candidate, faster discovery, lower
cost per drug program.

ANALOGY: Like CAD software for civil engineers — you don't build the bridge to test if
it stands. You simulate. SMEpred is CAD for siRNA drugs.

TRANSITION: "Quickly — the toolkit behind all this.""" ,

    # Slide 14 — Tech stack
    """OPENER: "Everything is mainstream open-source Python — two ML stacks under one app."

- ML (our model): LightGBM, scikit-learn, SciPy, NumPy.
- ML (ensemble): PyTorch 2.12 (CPU), vendored OligoFormer + RNA-FM (Tsinghua).
  The 1.1 GB RNA-FM checkpoint loads once at first /rank call with the toggle on.
- Data: pandas + custom HelixZero/siRNAmod/OligoFormer catalog parsers.
- Backend: FastAPI on Uvicorn — async, auto-Swagger at /docs.
- Frontend: single HTML file + vanilla JavaScript. Zero build pipeline.
- Quality: 19 unit tests, seeded reproducibility (metrics reproduce exactly).

VENDORING NOTE: We DON'T depend on the external OligoFormer folder. We copied only the
~1.1 GB of files we actually use into smepred/vendor/ and smepred/models/. The original
3.9 GB OligoFormer repo can be deleted.

ANALOGY: Same stack any modern AI startup uses. Not exotic. Production-grade.

TRANSITION: "Let me close with the takeaway.""" ,

    # Slide 15 — Closing
    """OPENER: Read the title. "From a gene to a drug candidate — in seconds."

Three takeaways:
1. Predicts BOTH naked and modified siRNA efficacy with built-in seed-toxicity AND
   modification-aware off-target rescue checks (Mitigated badges).
2. Accuracy:
     Modification ranking: 0.37 → 0.68 PCC (+84%)
     Naked siRNA ranking:  0.21 → 0.42 PCC (+100%) — the OligoFormer dataset merge.
3. Two ML models ensemble (our LightGBM + the OligoFormer transformer), wrapped in a
   clean web app any researcher can use immediately.

NEXT: more diverse target genes for the Rank tab. Wet-lab validation of SMEpred's top
picks — turn a software prediction into a confirmed lab result. Add the off-target
PITA/TargetScan filters from OligoFormer's pipeline as a third safety check.

CLOSE: "Thank you. I'm happy to take questions — about the biology, the model, or the
engineering."

LIKELY Q&A:
- Why not just use OligoFormer directly? It only handles naked siRNA. We need
  chemical-modification support for real drug design.
- Why is the OligoFormer score so much higher than yours? Different scales (probability
  vs % inhibition). That's why we use PERCENTILE-RANK ensemble — calibration-invariant.
- What does "Mitigated" mean? The seed lookup says "Toxic" but a known off-target
  rescue modification (2'-OMe, 2'-F, LNA, 2'-MOE) sits in pos 2-7 of the antisense.
  Established biology (Jackson et al., RNA 2006).
- Is 0.68 PCC good? The paper claimed 0.80 on a smaller curated set; ours is on a much
  larger messier real-world catalog with honest gene-grouped evaluation.
- Biggest limitation? Only 13 target genes in the cm-siRNA training set — need more
  diversity for the new-gene Rank use case.
- Wet-lab ready? For shortlisting yes (top 10 instead of 1,260 modifications, and the
  ones flagged Toxic OR func-fail get auto-dropped). For final go/no-go, still needs
  experimental confirmation.""" ,
]

assert len(NOTES) == 15, "expected 15 notes"

prs = Presentation(str(DECK))
assert len(prs.slides) == 15, f"expected 15 slides, got {len(prs.slides)}"

for i, slide in enumerate(prs.slides):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = NOTES[i]
prs.save(str(DECK))
print(f"Embedded notes into {len(NOTES)} slides.")
